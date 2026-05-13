"""
A0 dikey bitirme projesi posteri uretici.
Cikti: poster.pptx (proje koku)

A0 dikey: 84.1 cm x 118.9 cm
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Renk paleti ---
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
TEAL = RGBColor(0x06, 0xB6, 0xD4)
BG = RGBColor(0xF8, 0xFA, 0xFC)
TEXT = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
GREEN = RGBColor(0x10, 0xB9, 0x81)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_NAVY = RGBColor(0xE0, 0xE7, 0xFF)
LIGHT_TEAL = RGBColor(0xCF, 0xFA, 0xFE)


def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.1)
    tf.margin_bottom = Cm(0.1)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_box(slide, x, y, w, h, *, fill=WHITE, line=NAVY, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    box.shadow.inherit = False
    if box.has_text_frame:
        box.text_frame.text = ""
    return box


def add_filled_box_with_text(slide, x, y, w, h, text, *, fill=NAVY, text_color=WHITE,
                              size=22, bold=True, align=PP_ALIGN.CENTER):
    box = add_box(slide, x, y, w, h, fill=fill, line=None)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.2)
    tf.margin_right = Cm(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return box


def add_arrow_down(slide, x, y, length=Cm(1.0)):
    """Dikey ok cizer."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x - Cm(0.4), y, Cm(0.8), length)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TEAL
    arrow.line.fill.background()
    return arrow


def section_header(slide, x, y, w, text):
    """Sectiona baslik + alt cizgi."""
    add_text(slide, x, y, w, Cm(1.6), text,
             size=36, bold=True, color=NAVY)
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y + Cm(1.7), x + w, y + Cm(1.7))
    line.line.color.rgb = TEAL
    line.line.width = Pt(3)
    return y + Cm(2.2)


def main():
    prs = Presentation()
    # A0 dikey
    prs.slide_width = Cm(84.1)
    prs.slide_height = Cm(118.9)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # --- Arka plan ---
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    margin = Cm(2.5)
    inner_w = prs.slide_width - 2 * margin
    cursor = Cm(1.5)

    # ============== HEADER ==============
    header_h = Cm(11)
    header = add_box(slide, margin, cursor, inner_w, header_h, fill=NAVY, line=None)

    # Sol ust: universite placeholder
    add_text(slide, margin + Cm(0.8), cursor + Cm(0.6), Cm(20), Cm(1.2),
             "AMASYA UNIVERSITESI",
             size=22, bold=True, color=WHITE)
    add_text(slide, margin + Cm(0.8), cursor + Cm(2.0), Cm(30), Cm(1.0),
             "Muhendislik Fakultesi  |  Bilgisayar Muhendisligi",
             size=18, color=LIGHT_TEAL)

    # Ana baslik
    add_text(slide, margin + Cm(0.8), cursor + Cm(3.6), inner_w - Cm(1.6), Cm(4.2),
             "COK ORGANIZASYONLU AI DESTEKLI\nPROJE YONETIM SISTEMI",
             size=58, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Alt: ogrenci + danisman
    add_text(slide, margin + Cm(0.8), cursor + Cm(8.6), Cm(40), Cm(1.0),
             "Hazirlayan: Metin Karyagdi",
             size=22, bold=True, color=WHITE)
    add_text(slide, margin + Cm(0.8), cursor + Cm(9.8), Cm(40), Cm(1.0),
             "Danisman: [DANISMAN ADI]",
             size=20, color=LIGHT_TEAL)

    # Sag: yil + bolum kodu
    add_text(slide, margin + inner_w - Cm(15), cursor + Cm(8.6), Cm(14), Cm(1.0),
             "Bitirme Projesi  |  2026",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    cursor += header_h + Cm(1.0)

    # ============== OZET ==============
    cursor = section_header(slide, margin, cursor, inner_w, "OZET")
    ozet_box = add_box(slide, margin, cursor, inner_w, Cm(6), fill=LIGHT_NAVY, line=NAVY)
    add_text(slide, margin + Cm(0.5), cursor + Cm(0.4), inner_w - Cm(1.0), Cm(5.2),
             "Bu proje; cok-organizasyonlu (multi-tenant) bir proje yonetim sistemi olup, "
             "ekiplerin issue ve sprint takibini kolaylastirmak icin yerel calisan, "
             "ince-ayarli (fine-tuned) bir buyuk dil modeli ile entegre edilmistir. "
             "Mikroservis mimarisi, olay-tabanli (event-driven) iletisim ve organizasyon "
             "bazli izolasyon prensipleri uzerine kurulmustur. Veri gizliligi icin AI "
             "tamamen yerel cihazda calisir; hicbir ekip verisi disariya cikmaz.",
             size=22, color=TEXT, anchor=MSO_ANCHOR.TOP)
    cursor += Cm(7.0)

    # ============== PROBLEM + AMAC (2 sutun) ==============
    cursor = section_header(slide, margin, cursor, inner_w, "PROBLEM & AMAC")
    col_w = (inner_w - Cm(1.0)) / 2
    col_h = Cm(11)

    # Sol: Problem
    add_box(slide, margin, cursor, col_w, col_h, fill=WHITE, line=NAVY)
    add_text(slide, margin + Cm(0.5), cursor + Cm(0.3), col_w - Cm(1), Cm(1.2),
             "PROBLEM", size=26, bold=True, color=NAVY)
    add_text(slide, margin + Cm(0.5), cursor + Cm(1.7), col_w - Cm(1), Cm(9.0),
             "• Klasik proje yonetim araclari tek-tenant veya pahali SaaS\n"
             "• AI ozellikleri bulut modellerine bagli (OpenAI vs.) -> veri gizliligi riski\n"
             "• Ekipler farkli organizasyonlarda izole calisamiyor\n"
             "• Sprint analizi ve issue zenginlestirme manuel yapiliyor\n"
             "• Olay tabanli mimariler kucuk takimlar icin karmasik",
             size=20, color=TEXT)

    # Sag: Amac
    x2 = margin + col_w + Cm(1.0)
    add_box(slide, x2, cursor, col_w, col_h, fill=WHITE, line=TEAL)
    add_text(slide, x2 + Cm(0.5), cursor + Cm(0.3), col_w - Cm(1), Cm(1.2),
             "AMAC & KATKI", size=26, bold=True, color=TEAL)
    amac_lines = [
        "✓ Cok organizasyonlu, izole ekip yonetimi",
        "✓ Yerel calisan AI asistan (bp-agent)",
        "✓ Mikroservis + RabbitMQ olay-tabanli mimari",
        "✓ Otomatik sprint analizi ve issue zenginlestirme",
        "✓ Davet-tabanli rol yonetimi (Owner/Manager/Member)",
        "✓ SignalR ile gercek-zamanli bildirim",
    ]
    add_text(slide, x2 + Cm(0.5), cursor + Cm(1.7), col_w - Cm(1), Cm(9.0),
             "\n".join(amac_lines), size=20, color=TEXT)

    cursor += col_h + Cm(1.5)

    # ============== SISTEM MIMARISI (HERO) ==============
    cursor = section_header(slide, margin, cursor, inner_w, "SISTEM MIMARISI")
    arch_h = Cm(34)
    arch_box = add_box(slide, margin, cursor, inner_w, arch_h, fill=WHITE, line=NAVY, line_w=2)

    ax = margin + Cm(1.0)
    aw = inner_w - Cm(2.0)
    ay = cursor + Cm(0.8)

    # Layer 1: Frontend
    fe_h = Cm(2.6)
    add_filled_box_with_text(slide, ax, ay, aw, fe_h,
                              "FRONTEND  -  React + TypeScript + SignalR Client",
                              fill=TEAL, size=24)
    add_arrow_down(slide, ax + aw / 2, ay + fe_h + Cm(0.1), length=Cm(1.0))

    # Layer 2: Gateway
    gw_y = ay + fe_h + Cm(1.3)
    add_filled_box_with_text(slide, ax, gw_y, aw, fe_h,
                              "API GATEWAY  -  YARP Reverse Proxy + JWT Auth",
                              fill=NAVY, size=24)
    add_arrow_down(slide, ax + aw / 2, gw_y + fe_h + Cm(0.1), length=Cm(1.0))

    # Layer 3: Microservices (7 hucre)
    ms_y = gw_y + fe_h + Cm(1.3)
    ms_h = Cm(4.5)
    services = ["Identity", "Projects", "Issues", "Sprints", "Storage", "Notifications", "BFF"]
    cell_w = (aw - Cm(0.6) * (len(services) - 1)) / len(services)
    for i, name in enumerate(services):
        cx = ax + i * (cell_w + Cm(0.6))
        b = add_box(slide, cx, ms_y, cell_w, ms_h, fill=LIGHT_NAVY, line=NAVY, line_w=1.5)
        add_text(slide, cx, ms_y + Cm(0.4), cell_w, Cm(1.2),
                 name, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, cx, ms_y + Cm(1.6), cell_w, Cm(2.7),
                 ".NET 8\nWebAPI", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # MS altinda: "Mikroservisler (.NET 8)" etiketi
    add_text(slide, ax, ms_y + ms_h + Cm(0.1), aw, Cm(0.8),
             "Mikroservisler  -  .NET 8 + EF Core + CQRS/MediatR",
             size=16, color=MUTED, align=PP_ALIGN.CENTER)

    # Event bus
    eb_y = ms_y + ms_h + Cm(1.3)
    add_filled_box_with_text(slide, ax, eb_y, aw, Cm(2.4),
                              "RabbitMQ  -  Outbox Pattern + Event-Driven Communication",
                              fill=GREEN, size=22)
    add_arrow_down(slide, ax + aw / 2, eb_y + Cm(2.5), length=Cm(1.0))

    # AI Service (one cikan)
    ai_y = eb_y + Cm(3.7)
    ai_h = Cm(3.2)
    ai_box = add_box(slide, ax, ai_y, aw, ai_h, fill=NAVY, line=TEAL, line_w=3)
    tf = ai_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "AI SERVICE - bp-agent"
    r1.font.size = Pt(28); r1.font.bold = True; r1.font.color.rgb = TEAL; r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Ollama + Gemma3:4b Fine-tuned (LoRA + GGUF)"
    r2.font.size = Pt(20); r2.font.color.rgb = WHITE; r2.font.name = "Calibri"

    # Data layer (PostgreSQL + Redis)
    dl_y = ai_y + ai_h + Cm(1.0)
    dl_h = Cm(2.6)
    half = (aw - Cm(0.8)) / 2
    add_filled_box_with_text(slide, ax, dl_y, half, dl_h,
                              "PostgreSQL  -  Persistent Data",
                              fill=MUTED, size=22)
    add_filled_box_with_text(slide, ax + half + Cm(0.8), dl_y, half, dl_h,
                              "Redis  -  Cache & Session",
                              fill=MUTED, size=22)

    cursor += arch_h + Cm(1.5)

    # ============== AI ASISTAN DEMO (sohbet baloncuklari) ==============
    cursor = section_header(slide, margin, cursor, inner_w, "AI ASISTAN ORNEK DIYALOG")

    chat_h = Cm(13)
    add_box(slide, margin, cursor, inner_w, chat_h, fill=WHITE, line=NAVY)

    # Kullanici balonu (sol)
    u_box = add_box(slide, margin + Cm(1), cursor + Cm(0.8), Cm(34), Cm(4.5),
                    fill=LIGHT_NAVY, line=NAVY)
    add_text(slide, margin + Cm(1.4), cursor + Cm(1.0), Cm(33), Cm(1.0),
             "KULLANICI", size=18, bold=True, color=NAVY)
    add_text(slide, margin + Cm(1.4), cursor + Cm(2.0), Cm(33), Cm(3.0),
             "Sprint 12'nin durumunu ozetler misin? Hangi issue'lar bekliyor?",
             size=22, color=TEXT)

    # AI balonu (sag)
    a_box = add_box(slide, margin + inner_w - Cm(45), cursor + Cm(6.0), Cm(44), Cm(6.5),
                    fill=NAVY, line=TEAL, line_w=2)
    add_text(slide, margin + inner_w - Cm(44.5), cursor + Cm(6.2), Cm(43), Cm(1.0),
             "bp-agent (AI)", size=18, bold=True, color=TEAL)
    add_text(slide, margin + inner_w - Cm(44.5), cursor + Cm(7.2), Cm(43), Cm(5.0),
             "Sprint 12: 18 issue / 14 tamamlandi, 4 acik.\n"
             "  - 2 issue 'Blocked' (PROJ-142, PROJ-156)\n"
             "  - 1 issue 'In Progress' (Ahmet, gecikme: 2 gun)\n"
             "  - Tahmini bitirme: 9 Mayis (sprint sonu: 10 Mayis)\n"
             "Oneriler: Blocker'lari oncelendirin; sprint riski DUSUK.",
             size=20, color=WHITE)

    cursor += chat_h + Cm(1.5)

    # ============== TEKNOLOJI STACK ==============
    cursor = section_header(slide, margin, cursor, inner_w, "TEKNOLOJI STACK")
    stack_h = Cm(7)
    add_box(slide, margin, cursor, inner_w, stack_h, fill=WHITE, line=NAVY)

    techs = [
        ("Frontend", "React 18, TypeScript, Vite, TailwindCSS, SignalR Client"),
        ("Backend",  ".NET 8, ASP.NET Core, EF Core, MediatR (CQRS), AutoMapper"),
        ("AI",       "Python, FastAPI, Ollama, Gemma3:4b, LoRA Fine-tune, llama.cpp"),
        ("Altyapi",  "PostgreSQL, Redis, RabbitMQ, Docker, Seq (logging), MailHog"),
    ]
    row_h = (stack_h - Cm(0.8)) / len(techs)
    for i, (k, v) in enumerate(techs):
        ry = cursor + Cm(0.4) + i * row_h
        add_text(slide, margin + Cm(0.6), ry, Cm(8), row_h,
                 k, size=22, bold=True, color=TEAL,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, margin + Cm(9), ry, inner_w - Cm(9.5), row_h,
                 v, size=20, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

    cursor += stack_h + Cm(1.5)

    # ============== EKRAN GORUNTULERI ==============
    cursor = section_header(slide, margin, cursor, inner_w, "UYGULAMA EKRANLARI")
    ss_h = Cm(11)
    ss_w = (inner_w - Cm(2.0)) / 3
    placeholders = [
        ("Issue Board", "Sprint icindeki issue'larin\nKanban tahtasi"),
        ("AI Asistan", "bp-agent ile sohbet\nve oneriler"),
        ("Sprint Analytics", "Sprint metrikleri\nve burndown"),
    ]
    for i, (title, desc) in enumerate(placeholders):
        sx = margin + i * (ss_w + Cm(1.0))
        b = add_box(slide, sx, cursor, ss_w, ss_h, fill=LIGHT_TEAL, line=NAVY, line_w=1.5)
        add_text(slide, sx, cursor + Cm(2.5), ss_w, Cm(1.5),
                 "[ EKRAN GORUNTUSU ]", size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, sx, cursor + Cm(5.5), ss_w, Cm(1.4),
                 title, size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, sx, cursor + Cm(7.0), ss_w, Cm(3.0),
                 desc, size=18, color=TEXT, align=PP_ALIGN.CENTER)

    cursor += ss_h + Cm(1.5)

    # ============== SONUCLAR + QR ==============
    cursor = section_header(slide, margin, cursor, inner_w, "SONUCLAR")
    res_h = Cm(8)

    res_w = inner_w - Cm(12)
    add_box(slide, margin, cursor, res_w, res_h, fill=LIGHT_NAVY, line=NAVY)
    sonuc = [
        "✓  210 / 210 birim & entegrasyon testi gecti",
        "✓  7 mikroservis + Gateway + AI servisi calisir durumda",
        "✓  Yerel AI ile veri gizliligi (sifir disari cikan veri)",
        "✓  Cok-organizasyon izolasyonu (multi-tenant) saglandi",
        "✓  Olay-tabanli mimari (Outbox + RabbitMQ) entegre",
    ]
    add_text(slide, margin + Cm(0.6), cursor + Cm(0.4), res_w - Cm(1.2), res_h - Cm(0.8),
             "\n".join(sonuc), size=22, color=TEXT)

    # QR kutusu (placeholder)
    qx = margin + res_w + Cm(1.0)
    qw = inner_w - res_w - Cm(1.0)
    add_box(slide, qx, cursor, qw, res_h, fill=WHITE, line=TEAL, line_w=2)
    add_text(slide, qx, cursor + Cm(0.4), qw, Cm(1.2),
             "Proje Repo", size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, qx, cursor + Cm(1.6), qw, Cm(1.0),
             "[ QR KOD ]", size=16, color=MUTED, align=PP_ALIGN.CENTER)
    # QR yer tutucusu (kare)
    qr_size = Cm(5.5)
    qr_x = qx + (qw - qr_size) / 2
    qr_y = cursor + Cm(2.4)
    add_box(slide, qr_x, qr_y, qr_size, qr_size, fill=BG, line=MUTED, line_w=1)
    add_text(slide, qr_x, qr_y + qr_size / 2 - Cm(0.5), qr_size, Cm(1.0),
             "QR\nburaya", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # --- Kaydet ---
    out = "poster.pptx"
    prs.save(out)
    print(f"OK -> {out}")
    print(f"Boyut: {prs.slide_width / 360000:.1f} cm x {prs.slide_height / 360000:.1f} cm")


if __name__ == "__main__":
    main()
